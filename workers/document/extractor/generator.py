import io
import markdown
from fpdf import FPDF
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
import logging

logger = logging.getLogger("document-generator")

def generate_pdf(text: str, title: str) -> bytes:
    """Generate PDF from markdown text"""
    class PDF(FPDF):
        def header(self):
            self.set_font("helvetica", "B", 15)
            self.cell(0, 10, title, 0, 1, "C")
            self.ln(5)

    pdf = PDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=11)
    
    # Strip basic markdown for PDF (simplified)
    # A robust solution uses markdown/html2pdf, but for text:
    lines = text.split('\n')
    for line in lines:
        cleaned = line.replace('**', '').replace('__', '').replace('*', '')
        if cleaned.startswith('#'):
            pdf.set_font("helvetica", "B", 14)
            pdf.multi_cell(0, 10, cleaned.lstrip('# ').strip())
            pdf.set_font("helvetica", size=11)
        else:
            pdf.multi_cell(0, 6, cleaned)
            
    buf = io.BytesIO()
    pdf.output(buf)
    return buf.getvalue()

def generate_docx(text: str, title: str) -> bytes:
    """Generate DOCX from markdown text"""
    doc = Document()
    doc.add_heading(title, 0)
    
    lines = text.split('\n')
    for line in lines:
        cleaned = line.replace('**', '').replace('__', '')
        if cleaned.startswith('#'):
            level = min(cleaned.count('#'), 9)
            doc.add_heading(cleaned.lstrip('# ').strip(), level=level)
        elif cleaned.startswith('- '):
            doc.add_paragraph(cleaned[2:], style='List Bullet')
        else:
            if cleaned.strip():
                doc.add_paragraph(cleaned)
                
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()

def generate_xlsx(data: list[list[str]], sheet_name: str = "Sheet1") -> bytes:
    """Generate XLSX from a 2D array of data"""
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    
    for row in data:
        ws.append(row)
        
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()

def generate_pptx(slides_data: list[dict]) -> bytes:
    """Generate PPTX from a list of slides
       slides_data: [{'title': '...', 'content': ['point 1', 'point 2']}]
    """
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # Generate an initial title slide if it's the first one
    if slides_data and "main_title" in slides_data[0]:
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = slides_data[0].get("main_title", "Presentation")
        subtitle.text = slides_data[0].get("subtitle", "")
        slides_data = slides_data[1:]

    for slide_info in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_info.get("title", "Untitled Slide")
        tf = body_shape.text_frame
        
        for idx, point in enumerate(slide_info.get("content", [])):
            if idx == 0:
                tf.text = str(point)
            else:
                p = tf.add_paragraph()
                p.text = str(point)
                p.level = 0
                
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
