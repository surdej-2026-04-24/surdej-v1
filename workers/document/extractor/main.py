"""
Document Extraction Microservice — Office document → Markdown conversion.

Handles .docx, .xlsx, and .pptx files using Python libraries:
  - python-docx  → .docx paragraphs, tables, headings, styles
  - openpyxl     → .xlsx sheets, cells, merged ranges
  - python-pptx  → .pptx slides, text frames, tables, notes

Exposes HTTP endpoints called by the Node.js document worker.
Follows the same pattern as the pdf-extractor service.
"""

import io
import os
import logging
import tempfile
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse, Response
from pydantic import BaseModel
from typing import Optional, List
import uvicorn

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger("document-extractor")

app = FastAPI(title="Document Extractor", version="1.0.0")

# ─── Supported MIME types ─────────────────────────────────────

DOCX_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
XLSX_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
}
PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}
ALL_MIMES = DOCX_MIMES | XLSX_MIMES | PPTX_MIMES

EXTENSION_MAP = {
    ".docx": "docx",
    ".doc": "docx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".pptx": "pptx",
    ".ppt": "pptx",
}


@app.get("/health")
async def health():
    return {"status": "ok", "engine": "document-extractor", "formats": ["docx", "xlsx", "pptx"]}


@app.post("/extract")
async def extract_document(file: UploadFile = File(...)):
    """
    Extract markdown from an uploaded Office document.

    Accepts .docx, .xlsx, .pptx files.
    Returns JSON with markdown text, page/sheet/slide count, and quality score.
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="Filename required")

    ext = Path(file.filename).suffix.lower()
    doc_type = EXTENSION_MAP.get(ext)

    if not doc_type:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {ext}. Supported: {', '.join(EXTENSION_MAP.keys())}",
        )

    content = await file.read()
    if len(content) == 0:
        raise HTTPException(status_code=400, detail="Empty file")

    logger.info(f"Extracting {doc_type}: {file.filename} ({len(content)} bytes)")

    try:
        if doc_type == "docx":
            result = _extract_docx(content, file.filename)
        elif doc_type == "xlsx":
            result = _extract_xlsx(content, file.filename)
        elif doc_type == "pptx":
            result = _extract_pptx(content, file.filename)
        else:
            raise HTTPException(status_code=400, detail=f"Unknown type: {doc_type}")

        quality = _score_quality(result["markdown"])
        result["quality"] = quality

        logger.info(
            f"Extracted: {file.filename} → {len(result['markdown'])} chars, "
            f"quality={quality}, type={doc_type}"
        )

        return JSONResponse(content=result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Extraction failed for {file.filename}: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Extraction failed: {str(e)}")


class GenerateRequest(BaseModel):
    format: str # pdf, docx, xlsx, pptx
    title: str = "Document"
    text: Optional[str] = None
    data: Optional[List[List[str]]] = None
    slides: Optional[List[dict]] = None

@app.post("/generate")
async def generate_document_endpoint(req: GenerateRequest):
    import generator
    
    logger.info(f"Generating {req.format} document: {req.title}")
    
    try:
        if req.format == "pdf":
            content = generator.generate_pdf(req.text or "", req.title)
            return Response(content=content, media_type="application/pdf")
        elif req.format == "docx":
            content = generator.generate_docx(req.text or "", req.title)
            return Response(content=content, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        elif req.format == "xlsx":
            content = generator.generate_xlsx(req.data or [], req.title)
            return Response(content=content, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        elif req.format == "pptx":
            content = generator.generate_pptx(req.slides or [])
            return Response(content=content, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported format: {req.format}")
    except Exception as e:
        logger.error(f"Generation failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Generation failed: {str(e)}")


# ─── DOCX Extraction ─────────────────────────────────────────


def _extract_docx(content: bytes, filename: str) -> dict:
    """Extract markdown from a .docx file using python-docx, preserving formatting."""
    from docx import Document
    from docx.oxml.ns import qn
    import re

    doc = Document(io.BytesIO(content))
    lines: list[str] = []
    heading = Path(filename).stem
    has_heading = False

    # Build a unified iteration over body elements (paragraphs + tables in order)
    body = doc.element.body
    for child in body:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag

        if tag == 'p':
            # Find the matching paragraph object
            para = None
            for p in doc.paragraphs:
                if p._element is child:
                    para = p
                    break
            if para is None:
                continue

            text = para.text.strip()

            # Page break detection
            if any(br.get(qn('w:type')) == 'page' for br in child.findall('.//'+qn('w:br'))):
                lines.append('\n---\n')

            if not text:
                lines.append("")
                continue

            style_name = (para.style.name or "").lower() if para.style else ""

            # Heading detection
            if style_name.startswith("heading"):
                has_heading = True
                try:
                    level = int(style_name.replace("heading", "").strip())
                except ValueError:
                    level = 1
                level = min(level, 6)
                lines.append(f"{'#' * level} {text}")
                continue

            if style_name == "title":
                has_heading = True
                lines.append(f"# {text}")
                continue

            if style_name == "subtitle":
                lines.append(f"## {text}")
                continue

            # List detection with indent level
            numPr = child.find('.//' + qn('w:numPr'))
            if numPr is not None or style_name.startswith("list"):
                indent = 0
                ilvl = numPr.find(qn('w:ilvl')) if numPr is not None else None
                if ilvl is not None:
                    try:
                        indent = int(ilvl.get(qn('w:val'), '0'))
                    except (ValueError, TypeError):
                        indent = 0
                prefix = "    " * indent
                if style_name.startswith("list number") or "number" in style_name:
                    lines.append(f"{prefix}1. {_format_runs(para)}")
                else:
                    lines.append(f"{prefix}- {_format_runs(para)}")
                continue

            # Quote / block quote
            if style_name in ("quote", "intense quote", "block text"):
                lines.append(f"> {_format_runs(para)}")
                continue

            # Regular paragraph with inline formatting
            formatted = _format_runs(para)
            lines.append(formatted)

        elif tag == 'tbl':
            # Find the matching table object
            tbl = None
            for t in doc.tables:
                if t._element is child:
                    tbl = t
                    break
            if tbl:
                lines.append("")
                lines.append(_table_to_markdown(tbl))
                lines.append("")

    markdown = "\n".join(lines).strip()

    # Ensure document starts with a heading
    if not has_heading and markdown:
        markdown = f"# {heading}\n\n{markdown}"

    # Collapse excessive blank lines
    markdown = re.sub(r"\n{3,}", "\n\n", markdown)

    return {
        "markdown": markdown,
        "pageCount": len(doc.sections),
        "engine": "python-docx",
        "format": "docx",
    }


def _format_runs(para) -> str:
    """Convert paragraph runs to markdown with bold, italic, underline, strikethrough, and hyperlinks."""
    from docx.oxml.ns import qn

    parts: list[str] = []

    # Build a map of hyperlink elements to their URL
    rels = para.part.rels if hasattr(para, 'part') else {}

    for child in para._element:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag

        if tag == 'hyperlink':
            # Extract hyperlink URL
            r_id = child.get(qn('r:id'))
            url = ''
            if r_id and r_id in rels:
                url = rels[r_id].target_ref
            link_text_parts = []
            for run_elem in child.findall(qn('w:r')):
                t = run_elem.findtext(qn('w:t'), default='')
                link_text_parts.append(t)
            link_text = ''.join(link_text_parts)
            if url and link_text:
                parts.append(f"[{link_text}]({url})")
            elif link_text:
                parts.append(link_text)

        elif tag == 'r':
            # Regular run
            from docx.text.run import Run
            run = Run(child, para)
            text = run.text
            if not text:
                continue

            is_bold = run.bold
            is_italic = run.italic
            is_underline = run.underline
            is_strike = run.font.strike if run.font else False

            if is_bold and is_italic:
                text = f"***{text}***"
            elif is_bold:
                text = f"**{text}**"
            elif is_italic:
                text = f"*{text}*"

            if is_strike:
                text = f"~~{text}~~"
            if is_underline and not is_bold:
                text = f"<u>{text}</u>"

            parts.append(text)

    return "".join(parts) or para.text


def _table_to_markdown(table) -> str:
    """Convert a docx table to markdown table."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    # Build markdown table
    col_count = max(len(r) for r in rows)
    lines: list[str] = []

    # Header row
    header = rows[0] if rows else [""] * col_count
    header += [""] * (col_count - len(header))
    lines.append("| " + " | ".join(h.replace("|", "\\|") for h in header) + " |")
    lines.append("| " + " | ".join(["---"] * col_count) + " |")

    # Data rows
    for row in rows[1:]:
        row += [""] * (col_count - len(row))
        lines.append("| " + " | ".join(c.replace("|", "\\|") for c in row) + " |")

    return "\n".join(lines)


# ─── XLSX Extraction ─────────────────────────────────────────


def _extract_xlsx(content: bytes, filename: str) -> dict:
    """Extract markdown from a .xlsx file using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines: list[str] = []
    heading = Path(filename).stem
    lines.append(f"# {heading}")
    lines.append("")

    sheet_count = len(wb.sheetnames)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## {sheet_name}")
        lines.append("")

        # Read all rows
        all_rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            # Skip completely empty rows
            if any(c.strip() for c in cells):
                all_rows.append(cells)

        if not all_rows:
            lines.append("*(Empty sheet)*")
            lines.append("")
            continue

        # Determine column count (max cells in any row)
        col_count = max(len(r) for r in all_rows)

        # Build markdown table
        # First row as header
        header = all_rows[0]
        header += [""] * (col_count - len(header))
        lines.append("| " + " | ".join(h.replace("|", "\\|").replace("\n", " ") for h in header) + " |")
        lines.append("| " + " | ".join(["---"] * col_count) + " |")

        for row in all_rows[1:]:
            row += [""] * (col_count - len(row))
            lines.append(
                "| "
                + " | ".join(c.replace("|", "\\|").replace("\n", " ") for c in row)
                + " |"
            )

        lines.append("")

    wb.close()

    markdown = "\n".join(lines).strip()

    return {
        "markdown": markdown,
        "sheetCount": sheet_count,
        "pageCount": sheet_count,
        "engine": "openpyxl",
        "format": "xlsx",
    }


# ─── PPTX Extraction ─────────────────────────────────────────


def _extract_pptx(content: bytes, filename: str) -> dict:
    """Extract markdown from a .pptx file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(io.BytesIO(content))
    lines: list[str] = []
    heading = Path(filename).stem
    lines.append(f"# {heading}")
    lines.append("")

    slide_count = len(prs.slides)

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {slide_num}")
        lines.append("")

        for shape in slide.shapes:
            # Text frames (titles, body text, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    # Detect title shape
                    if shape.shape_id == 0 or (hasattr(shape, "placeholder_format") and
                            shape.placeholder_format is not None and
                            shape.placeholder_format.idx in (0, 1)):
                        lines.append(f"### {text}")
                    else:
                        # Check bullet level
                        level = para.level if para.level else 0
                        indent = "  " * level
                        lines.append(f"{indent}- {text}" if level > 0 else text)

            # Tables
            if shape.has_table:
                table = shape.table
                rows: list[list[str]] = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows.append(cells)

                if rows:
                    col_count = max(len(r) for r in rows)
                    header = rows[0] + [""] * (col_count - len(rows[0]))
                    lines.append("")
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("| " + " | ".join(["---"] * col_count) + " |")
                    for row in rows[1:]:
                        row += [""] * (col_count - len(row))
                        lines.append("| " + " | ".join(row) + " |")
                    lines.append("")

        # Slide notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("")
                lines.append(f"> **Notes:** {notes_text}")

        lines.append("")

    markdown = "\n".join(lines).strip()

    # Collapse excessive blank lines
    import re
    markdown = re.sub(r"\n{3,}", "\n\n", markdown)

    return {
        "markdown": markdown,
        "slideCount": slide_count,
        "pageCount": slide_count,
        "engine": "python-pptx",
        "format": "pptx",
    }


# ─── Quality Scoring ─────────────────────────────────────────


def _score_quality(text: str) -> int:
    """Score markdown quality 0–100."""
    import re

    if not text or not text.strip():
        return 0

    words = text.split()
    unique_words = set(w.lower() for w in words)
    score = 0.0

    # Length factor (max 25)
    score += min(len(words) / 100, 1.0) * 25

    # Diversity (max 25)
    if words:
        score += (len(unique_words) / len(words)) * 25

    # Structure indicators (max 30)
    has_headings = text.count("\n#") > 0
    has_tables = "|" in text and "---" in text
    has_lists = any(line.strip().startswith(("- ", "* ", "1.")) for line in text.split("\n"))
    has_bold = "**" in text
    has_paragraphs = text.count("\n\n") > 2
    has_danish = any(c in text for c in "æøåÆØÅ")
    has_numbers = any(c.isdigit() for c in text)

    if has_headings: score += 6
    if has_tables: score += 6
    if has_lists: score += 4
    if has_bold: score += 3
    if has_paragraphs: score += 4
    if has_danish: score += 4
    if has_numbers: score += 3

    # Penalty for garbage
    garbage_chars = len(re.findall(r"[^\w\sæøåÆØÅ.,;:!?\-\[\]\"'/\\@#€$%&+*=<>|(){}~`^]", text))
    garbage_ratio = garbage_chars / max(len(text), 1)
    score -= garbage_ratio * 40

    # Penalty for very short text
    if len(words) < 20:
        score -= 20

    return max(0, min(100, round(score)))


if __name__ == "__main__":
    port = int(os.environ.get("EXTRACTOR_PORT", "8091"))
    logger.info(f"Starting Document Extractor on port {port}")
    uvicorn.run(app, host="0.0.0.0", port=port, log_level="info")
