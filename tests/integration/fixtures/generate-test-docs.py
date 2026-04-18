"""
Generate test fixture files for the document processing integration tests.

Creates .docx, .xlsx, and .pptx files with known content that the tests
can verify after extraction.

Run: python3 tests/integration/fixtures/generate-test-docs.py
"""

import os

FIXTURES_DIR = os.path.dirname(os.path.abspath(__file__))


def create_docx():
    from docx import Document

    doc = Document()
    doc.add_heading("Integration Test Report", 0)
    doc.add_heading("Summary", level=1)
    doc.add_paragraph("Budget: 142.857 DKK")
    doc.add_paragraph("Project: SURDEJ-DOCX-TEST-2026")
    doc.add_paragraph("This document tests the document extraction pipeline.")

    doc.add_heading("Team", level=2)
    table = doc.add_table(rows=3, cols=3)
    table.rows[0].cells[0].text = "Name"
    table.rows[0].cells[1].text = "Role"
    table.rows[0].cells[2].text = "Status"
    table.rows[1].cells[0].text = "Alice"
    table.rows[1].cells[1].text = "Developer"
    table.rows[1].cells[2].text = "Active"
    table.rows[2].cells[0].text = "Bob"
    table.rows[2].cells[1].text = "Designer"
    table.rows[2].cells[2].text = "Active"

    path = os.path.join(FIXTURES_DIR, "test-report.docx")
    doc.save(path)
    print(f"  ✓ {path}")


def create_xlsx():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Category", "Amount", "Currency"])
    ws.append(["Staff", "150000", "DKK"])
    ws.append(["Equipment", "42000", "DKK"])
    ws.append(["Software", "28000", "DKK"])
    ws.append(["Total", "220000", "DKK"])

    ws2 = wb.create_sheet("Projects")
    ws2.append(["Project", "Status", "Lead"])
    ws2.append(["SURDEJ-XLSX-TEST-2026", "Active", "Alice"])
    ws2.append(["Framework v2", "Planning", "Bob"])

    path = os.path.join(FIXTURES_DIR, "test-budget.xlsx")
    wb.save(path)
    print(f"  ✓ {path}")


def create_pptx():
    from pptx import Presentation

    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Project Overview"
    slide1.placeholders[1].text = (
        "Budget: 99.000 DKK\n"
        "Project: SURDEJ-PPTX-TEST-2026\n"
        "Status: Active"
    )

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Findings"
    slide2.placeholders[1].text = (
        "1. Performance improved 35%\n"
        "2. User satisfaction at 92%\n"
        "3. Cost reduced by 18%"
    )

    path = os.path.join(FIXTURES_DIR, "test-presentation.pptx")
    prs.save(path)
    print(f"  ✓ {path}")


if __name__ == "__main__":
    print("Generating test fixture documents...")
    create_docx()
    create_xlsx()
    create_pptx()
    print("Done.")
